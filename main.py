#!/usr/bin/env python

import argparse
import mistune
from mistune import escape
from pptx import Presentation

class DebugRenderer(object):
  def __init__(self, **kwargs):
    self.options = kwargs

  def placeholder(self):
    print "placeholder"
    return ''

  def block_code(self, code, lang=None):
    print "block", code, lang
    code = code.rstrip('\n')
    if not lang:
      code = escape(code, smart_amp=False)
      return '<pre><code>%s\n</code></pre>\n' % code
    code = escape(code, quote=True, smart_amp=False)
    return '<pre><code class="lang-%s">%s\n</code></pre>\n' % (lang, code)

  def block_quote(self, text):
    print "block_quote", text
    return '<blockquote>%s\n</blockquote>\n' % text.rstrip('\n')

  def block_html(self, html):
    print "block_html", html
    if self.options.get('skip_style') and \
       html.lower().startswith('<style'):
      return ''
    if self.options.get('escape'):
      return escape(html)
    return html

  def header(self, text, level, raw=None):
    print "header", text, level, raw
    return '<h%d>%s</h%d>\n' % (level, text, level)

  def hrule(self):
    print "hrule"
    if self.options.get('use_xhtml'):
      return '<hr />\n'
    return '<hr>\n'

  def list(self, body, ordered=True):
    print "list", body, ordered
    tag = 'ul'
    if ordered:
      tag = 'ol'
    return '<%s>\n%s</%s>\n' % (tag, body, tag)

  def list_item(self, text):
    print "list_item", text
    return '<li>%s</li>\n' % text

  def paragraph(self, text):
    print "paragraph", text
    return '<p>%s</p>\n' % text.strip(' ')

  def table(self, header, body):
    print "table", header, body
    return (
      '<table>\n<thead>%s</thead>\n'
      '<tbody>\n%s</tbody>\n</table>\n'
    ) % (header, body)

  def table_row(self, content):
    print "table_row", content
    return '<tr>\n%s</tr>\n' % content

  def table_cell(self, content, **flags):
    print "table_cell", content, flags
    if flags['header']:
      tag = 'th'
    else:
      tag = 'td'
    align = flags['align']
    if not align:
      return '<%s>%s</%s>\n' % (tag, content, tag)
    return '<%s style="text-align:%s">%s</%s>\n' % (
      tag, align, content, tag
    )

  def double_emphasis(self, text):
    print "double_emphasis", text
    return '<strong>%s</strong>' % text

  def emphasis(self, text):
    print "emphasis", text
    return '<em>%s</em>' % text

  def codespan(self, text):
    print "codespace", text
    text = escape(text.rstrip(), smart_amp=False)
    return '<code>%s</code>' % text

  def linebreak(self):
    print "linebreak"
    return '<br>\n'

  def strikethrough(self, text):
    print "strikethrough", text
    return '<del>%s</del>' % text

  def text(self, text):
    print "text", text
    return escape(text)

  def autolink(self, link, is_email=False):
    print "autolink", link, is_email
    text = link = escape(link)
    if is_email:
      link = 'mailto:%s' % link
    return '<a href="%s">%s</a>' % (link, text)

  def link(self, link, title, text):
    print "link", link, title, text
    link = escape_link(link, quote=True)
    if not title:
      return '<a href="%s">%s</a>' % (link, text)
    title = escape(title, quote=True)
    return '<a href="%s" title="%s">%s</a>' % (link, title, text)

  def image(self, src, title, text):
    print "image", src, title, text
    src = escape_link(src, quote=True)
    text = escape(text, quote=True)
    if title:
      title = escape(title, quote=True)
      html = '<img src="%s" alt="%s" title="%s"' % (src, text, title)
    else:
      html = '<img src="%s" alt="%s"' % (src, text)
    if self.options.get('use_xhtml'):
      return '%s />' % html
    return '%s>' % html

  def inline_html(self, html):
    print "inline_html", html
    if self.options.get('escape'):
      return escape(html)
    return html

  def newline(self):
    print "newline"
    return ''

  def footnote_ref(self, key, index):
    print "footnote_ref", key, index
    html = (
      '<sup class="footnote-ref" id="fnref-%s">'
      '<a href="#fn-%s" rel="footnote">%d</a></sup>'
    ) % (escape(key), escape(key), index)
    return html

  def footnote_item(self, key, text):
    print "footnote_item", key, text
    back = (
      '<a href="#fnref-%s" rev="footnote">&#8617;</a>'
    ) % escape(key)
    text = text.rstrip()
    if text.endswith('</p>'):
      text = re.sub(r'<\/p>$', r'%s</p>' % back, text)
    else:
      text = '%s<p>%s</p>' % (text, back)
    html = '<li id="fn-%s">%s</li>\n' % (escape(key), text)
    return html

  def footnotes(self, text):
    print "footnotes", text
    html = '<div class="footnotes">\n%s<ol>%s</ol>\n</div>\n'
    return html % (self.hrule(), text)


class Renderer(mistune.Renderer):
  def __init__(self, presentation, **kwargs):
    super(Renderer, self).__init__(**kwargs);
    self.presentation = presentation
    self.title_layout = self.presentation.slide_layouts[0]

  def placeholder(self):
    print "placeholder"
    return []

  def block_code(self, code, language=None):
    return []
  def block_quote(self, text):
    return []
  def block_html(self, html):
    print "block_html", html
    return [('text', html)]
  def hrule(self):
    return []
  def table(self, header, body):
    return []
  def table_row(self, content):
    return []
  def table_cell(self, content, **flags):
    return []

  def header(self, text, level, raw=None):
    print "header", text, level, raw
    # slide = self.presentation.slides.add_slide(self.title_layout)
    # title = slide.shapes.title
    # subtitle = slide.placeholders[1]
    # title.text = text
    return [('slide', text)]

  def list(self, body, ordered=True):
    print "list", body
    return [("list", body)]

  def list_item(self, text):
    list_item = []
    for item in text:
      if item[0] == 'text':
        if len(list_item) == 0:
          list_item.append(item)
        elif len(list_item) == 1:
          list_item[0] = ('text', list_item[0][1] + item[1])
        else:
          list_item.append(item)
      else:
        list_item.append(item)

    print "list_item", text, list_item
    return [("li", list_item)]

  def paragraph(self, text):
    return [text]
  def double_emphasis(self, text):
    print "double_emphasis", text
    return '<strong>%s</strong>' % text

  def emphasis(self, text):
    print "emphasis", text
    return '<em>%s</em>' % text

  def codespan(self, text):
    print "codespace", text
    text = escape(text.rstrip(), smart_amp=False)
    return '<code>%s</code>' % text

  def linebreak(self):
    print "linebreak"
    return '<br>\n'

  def strikethrough(self, text):
    print "strikethrough", text
    return '<del>%s</del>' % text

  def text(self, text):
    print "text", text
    return [('text', text)]

  def autolink(self, link, is_email=False):
    print "autolink", link, is_email
    text = link = escape(link)
    if is_email:
      link = 'mailto:%s' % link
    return '<a href="%s">%s</a>' % (link, text)

  def link(self, link, title, text):
    print "link", link, title, text
    link = escape_link(link, quote=True)
    if not title:
      return '<a href="%s">%s</a>' % (link, text)
    title = escape(title, quote=True)
    return '<a href="%s" title="%s">%s</a>' % (link, title, text)

  def inline_html(self, html):
    print "inline_html", html
    return html
    if self.options.get('escape'):
      return escape(html)
    return html

  def newline(self):
    print "newline"
    return ''

class DebugMarkdown(mistune.Markdown):
  def pop(self):
    print self.tokens
    return None

def build_presentation(parsed):
  presentation = Presentation()
  title_layout = presentation.slide_layouts[0]
  bullet_layout = presentation.slide_layouts[1]
  have_title = False
  last_slide = None
  for item in parsed:
    if item[0] == 'slide':
      if not have_title:
        slide = presentation.slides.add_slide(title_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = item[1][0][1]
        last_slide = slide
        have_title = True
      else:
        slide = presentation.slides.add_slide(bullet_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        last_slide = slide
        title.text = item[1][0][1]
    elif item[0] == 'list':
      from pprint import pprint
      pprint(item)
      body_shape = last_slide.shapes.placeholders[1]
      tf = body_shape.text_frame
      tf.text = ''
      def build_item(data, tf, level, paragraph):
        for datum in data:
          if datum[0] == 'text':
            if level == 0:
              tf.text = datum[1]
            else:
              paragraph.text = datum[1]
          elif datum[0] == 'list':
            build_bullets(datum[1], tf, level+1)
      def build_bullets(data, tf, level):
        for datum in data:
          if datum[0] == 'list':
            build_bullets(datum[1], tf, level+1)
          elif datum[0] == 'li':
            if level == 0:
              build_item(datum[1], tf, level, None)
            else:
              print "Adding Paragraph", datum, "->", data
              p = tf.add_paragraph()
              p.text = ''
              p.level = level
              build_item(datum[1], tf, level, p)
      build_bullets(item[1], tf, 0)

    else:
      print "ERRR: how do I handle", item
    '''

HI ('list', [('li', [('text', 'PPFE Traffic Management System'), ('list', [('li', [('text', 'Modern Front End'), ('list', [('li', [('text', 'HTTP/2 ( 2'), ('text', '<-> 1.1 tunnel
)')]), ('li', [('text', 'IPv6 (v6 '), ('text', '<-> v4 encapsulation)')]), ('li', [('text', 'TLSv1.2 (can upgrade or downgrade back end targets)')]), ('li', [('text', 'Cluster serv
es single IP via Anycast (Masterless clustering)')])])]), ('li', [('text', 'Firewall (L4 and L7)'), ('list', [('li', [('text', 'Global Rate Limiting (Millions of IPs in seconds)')]
), ('li', [('text', 'Deterministic Rule based firewall (Dynamic and configurable via ZK)')]), ('li', [('text', 'Behavioral Based firewall (Neural Network backed, continual deep lea
rning)')]), ('li', [('text', 'Cluster wide decisions')])])]), ('li', [('text', 'Load Balancer'), ('list', [('li', [('text', "Dynamic Async Health Checks (Protocol Specific ECV's i.
e OCC)")]), ('li', [('text', 'Predictive Load Balancing Algorithms')]), ('li', [('text', 'Dynamic Pool Configuration via Zookeeper')]), ('li', [('text', 'Weights, Soft Ramps and St
icky Sessions (and sticky failover)')])])])])])])
('list', [('li', [('text', 'PPFE Traffic Management System'), ('list', [('li', [('text', 'Modern Front End'), ('list', [('li', [('text', 'HTTP/2 ( 2'), ('text', '<-> 1.1 tunnel )')
]), ('li', [('text', 'IPv6 (v6 '), ('text', '<-> v4 encapsulation)')]), ('li', [('text', 'TLSv1.2 (can upgrade or downgrade back end targets)')]), ('li', [('text', 'Cluster serves
single IP via Anycast (Masterless clustering)')])])]), ('li', [('text', 'Firewall (L4 and L7)'), ('list', [('li', [('text', 'Global Rate Limiting (Millions of IPs in seconds)')]),
('li', [('text', 'Deterministic Rule based firewall (Dynamic and configurable via ZK)')]), ('li', [('text', 'Behavioral Based firewall (Neural Network backed, continual deep learni
ng)')]), ('li', [('text', 'Cluster wide decisions')])])]), ('li', [('text', 'Load Balancer'), ('list', [('li', [('text', "Dynamic Async Health Checks (Protocol Specific ECV's i.e O
CC)")]), ('li', [('text', 'Predictive Load Balancing Algorithms')]), ('li', [('text', 'Dynamic Pool Configuration via Zookeeper')]), ('li', [('text', 'Weights, Soft Ramps and Stick
y Sessions (and sticky failover)')])])])])])])


prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2'''
    print item
  return presentation

def process_markdown(md_file):
  presentation = Presentation()
  markdown = mistune.Markdown(renderer=Renderer(presentation))
  #markdown = DebugMarkdown(renderer=DebugRenderer(foo=presentation))
  result = markdown(open(md_file).read())
  #markdown(fuck)
  print "BUILDING"
  return build_presentation(result)

def parse_args():
  parser = argparse.ArgumentParser()
  parser.add_argument('md_file', type=str)
  parser.add_argument('pptx_file', type=str)
  parser.add_argument('-n', '--dry-run', action='store_const', const=True, default=False)

  return parser.parse_args()

def main():
  args = parse_args()
  process_markdown(args.md_file).save(args.pptx_file)

if __name__ == '__main__':
  main()

def poop():
  prs = Presentation()
  title_slide_layout = prs.slide_layouts[0]
  slide = prs.slides.add_slide(title_slide_layout)
  title = slide.shapes.title
  subtitle = slide.placeholders[1]

  title.text = "Hello, World!"
  subtitle.text = "python-pptx was here!"

  prs.save('test.pptx')
